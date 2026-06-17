"""Emit a JSON inventory of a template .pptx's sample slides (Step 2 fallback).

Used when `soffice`/`poppler` are unavailable for the vision path. Walks each
slide (in presentation order) with python-pptx and reports what objects it
holds, so the matching step can pick a template sample per slide entry using
text instead of an image.

Indices are 0-based presentation order, matching
`render_template_thumbnails.py` labels and the Step-4 build's mapping.json
`template_index`.

Usage:
    python inspect_template.py template.pptx        # JSON to stdout
    python inspect_template.py template.pptx --pretty
"""

import argparse
import json
import sys

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE, PP_PLACEHOLDER


def _ph_type(shape):
    try:
        if shape.is_placeholder:
            return str(shape.placeholder_format.type)
    except Exception:
        pass
    return None


def _is_title(shape) -> bool:
    try:
        if shape.is_placeholder:
            return shape.placeholder_format.type in (
                PP_PLACEHOLDER.TITLE,
                PP_PLACEHOLDER.CENTER_TITLE,
            )
    except Exception:
        pass
    return False


def _sample_text(shape, limit=80):
    try:
        if shape.has_text_frame:
            txt = " ".join(p.text for p in shape.text_frame.paragraphs).strip()
            txt = " ".join(txt.split())
            return txt[:limit]
    except Exception:
        pass
    return ""


def inspect_slide(slide, index):
    has_chart = has_table = has_picture = False
    n_text_ph = 0
    title = ""
    texts = []

    for shape in slide.shapes:
        if getattr(shape, "has_chart", False):
            has_chart = True
        if getattr(shape, "has_table", False):
            has_table = True
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            has_picture = True
        if _is_title(shape):
            title = _sample_text(shape)
        elif getattr(shape, "has_text_frame", False) and shape.text_frame.text.strip():
            n_text_ph += 1
            sample = _sample_text(shape)
            if sample:
                texts.append(sample)

    try:
        layout_name = slide.slide_layout.name
    except Exception:
        layout_name = None

    role = _guess_role(
        has_chart, has_table, has_picture, n_text_ph, title, layout_name
    )

    return {
        "index": index,
        "layout_name": layout_name,
        "title": title,
        "has_chart": has_chart,
        "has_table": has_table,
        "has_picture": has_picture,
        "n_text_placeholders": n_text_ph,
        "texts": texts,
        "suggested_role": role,
    }


def _guess_role(has_chart, has_table, has_picture, n_text_ph, title, layout_name):
    name = (layout_name or "").lower()
    if has_chart:
        return "chart"
    if has_table:
        return "table"
    if has_picture:
        return "image"
    if n_text_ph == 0:
        return "section_divider"
    if "section" in name or "divider" in name:
        return "section_divider"
    if "title" in name and "content" not in name and n_text_ph <= 1:
        return "title"
    return "bullets"


def main():
    parser = argparse.ArgumentParser(description="Inventory a template .pptx.")
    parser.add_argument("input", help="Template .pptx")
    parser.add_argument("--pretty", action="store_true", help="Indented JSON")
    args = parser.parse_args()

    try:
        prs = Presentation(args.input)
    except Exception as e:
        print(f"Error: cannot open {args.input}: {e}", file=sys.stderr)
        sys.exit(1)

    inventory = [inspect_slide(s, i) for i, s in enumerate(prs.slides)]
    print(json.dumps(inventory, indent=2 if args.pretty else None, ensure_ascii=False))


if __name__ == "__main__":
    main()
